"""Premium sample deck — demonstrates layout, typography, color system."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Design system -----------------------------------------------------------
INK       = RGBColor(0x0B, 0x1F, 0x3A)  # deep navy — primary text / backgrounds
PAPER     = RGBColor(0xFA, 0xF7, 0xF2)  # warm off-white — light backgrounds
MUTED     = RGBColor(0x6B, 0x73, 0x80)  # slate — secondary text
ACCENT    = RGBColor(0xC9, 0xA2, 0x6B)  # muted gold — accents, rules
ACCENT_2  = RGBColor(0xE0, 0x5B, 0x49)  # warm coral — highlight
HAIRLINE  = RGBColor(0xD9, 0xD2, 0xC7)  # subtle divider

FONT_DISPLAY = "Georgia"       # serif for display — premium feel
FONT_BODY    = "Segoe UI"      # clean sans for body
FONT_MONO    = "Consolas"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9

# --- Helpers -----------------------------------------------------------------
def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(s, color)
    return s

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=14, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
    return tb

def add_line(slide, x, y, w, h, color=ACCENT, weight=1.25):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def page_number(slide, n, total, color=MUTED):
    add_text(slide, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.3),
             f"{n:02d} / {total:02d}", font=FONT_BODY, size=9,
             color=color, align=PP_ALIGN.RIGHT, letter_spacing=200)

def footer(slide, label, color=MUTED):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             label, font=FONT_BODY, size=9, color=color, letter_spacing=300)

# --- Build -------------------------------------------------------------------
prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
blank = prs.slide_layouts[6]

TOTAL = 6

# ---------- Slide 1 — Cover --------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
# gold rule
add_rect(s, Inches(0.8), Inches(1.0), Inches(0.6), Emu(22860), ACCENT)
add_text(s, Inches(0.8), Inches(1.2), Inches(6), Inches(0.4),
         "AIBOS  ·  STRATEGY DECK", font=FONT_BODY, size=10,
         color=ACCENT, letter_spacing=600, bold=True)
add_text(s, Inches(0.8), Inches(2.4), Inches(11), Inches(2.2),
         "Designing the\nnext decade of work.",
         font=FONT_DISPLAY, size=60, color=PAPER)
add_text(s, Inches(0.8), Inches(5.2), Inches(10), Inches(0.5),
         "A premium sample built with python-pptx",
         font=FONT_BODY, size=14, color=HAIRLINE)
add_text(s, Inches(0.8), Inches(6.9), Inches(6), Inches(0.3),
         "April 2026   ·   Confidential",
         font=FONT_BODY, size=9, color=MUTED, letter_spacing=400)
add_text(s, Inches(11.3), Inches(6.9), Inches(2), Inches(0.3),
         "v1.0", font=FONT_BODY, size=9, color=MUTED,
         align=PP_ALIGN.RIGHT, letter_spacing=400)

# ---------- Slide 2 — Contents ----------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PAPER)
add_text(s, Inches(0.8), Inches(0.7), Inches(4), Inches(0.4),
         "CONTENTS", font=FONT_BODY, size=10, color=ACCENT,
         letter_spacing=600, bold=True)
add_line(s, Inches(0.8), Inches(1.15), Inches(0.5), Emu(15240), INK)
add_text(s, Inches(0.8), Inches(1.5), Inches(8), Inches(1.0),
         "What we'll cover", font=FONT_DISPLAY, size=40, color=INK)

items = [
    ("01", "Context",       "Where the market stands today"),
    ("02", "Opportunity",   "The gap we're positioned to close"),
    ("03", "Approach",      "Three principles shaping the work"),
    ("04", "Results",       "Early signal from pilot customers"),
]
y = Inches(3.1)
for num, title, sub in items:
    add_text(s, Inches(0.8), y, Inches(0.9), Inches(0.5),
             num, font=FONT_DISPLAY, size=22, color=ACCENT)
    add_text(s, Inches(1.9), y, Inches(5), Inches(0.5),
             title, font=FONT_BODY, size=16, color=INK, bold=True)
    add_text(s, Inches(5.2), y, Inches(7), Inches(0.5),
             sub, font=FONT_BODY, size=13, color=MUTED)
    add_line(s, Inches(0.8), y + Inches(0.75), Inches(11.7), Emu(6350), HAIRLINE)
    y += Inches(0.95)

footer(s, "CONTENTS")
page_number(s, 2, TOTAL)

# ---------- Slide 3 — Section divider ---------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
add_rect(s, 0, 0, Inches(4.5), SLIDE_H, RGBColor(0x08, 0x18, 0x2E))
add_text(s, Inches(0.8), Inches(1.0), Inches(3), Inches(0.4),
         "SECTION 02", font=FONT_BODY, size=10, color=ACCENT,
         letter_spacing=600, bold=True)
add_line(s, Inches(0.8), Inches(1.5), Inches(0.6), Emu(22860), ACCENT)
add_text(s, Inches(0.8), Inches(2.9), Inches(11), Inches(3),
         "Opportunity",
         font=FONT_DISPLAY, size=72, color=PAPER)
add_text(s, Inches(0.8), Inches(5.2), Inches(9), Inches(1.3),
         "A $48B category is being rebuilt from the data layer up —\nand the incumbents can't move fast enough to defend it.",
         font=FONT_DISPLAY, size=18, color=HAIRLINE)
page_number(s, 3, TOTAL, color=HAIRLINE)

# ---------- Slide 4 — Three-column content ----------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PAPER)
add_text(s, Inches(0.8), Inches(0.7), Inches(4), Inches(0.4),
         "APPROACH", font=FONT_BODY, size=10, color=ACCENT,
         letter_spacing=600, bold=True)
add_line(s, Inches(0.8), Inches(1.15), Inches(0.5), Emu(15240), INK)
add_text(s, Inches(0.8), Inches(1.5), Inches(10), Inches(1.0),
         "Three principles shaping the work",
         font=FONT_DISPLAY, size=34, color=INK)

cols = [
    ("Clarity",
     "Every interface answers one question: what should I do next?",
     "We strip decisions down to the one that matters in context — no dashboards for their own sake."),
    ("Trust",
     "Every number is reproducible from source to screen in under three clicks.",
     "Auditability isn't a compliance checkbox — it's how users develop confidence in the product."),
    ("Speed",
     "Every interaction budget is measured in milliseconds, not seconds.",
     "Latency is a feature. We instrument p95 on the hot paths and treat regressions as bugs."),
]
col_w = Inches(3.8)
gap   = Inches(0.35)
x0    = Inches(0.8)
y0    = Inches(3.1)
for i, (h, lead, body) in enumerate(cols):
    x = x0 + i * (col_w + gap)
    add_rect(s, x, y0, Inches(0.35), Emu(38100), ACCENT)
    add_text(s, x, y0 + Inches(0.25), col_w, Inches(0.5),
             f"0{i+1}", font=FONT_BODY, size=10, color=ACCENT,
             letter_spacing=400, bold=True)
    add_text(s, x, y0 + Inches(0.55), col_w, Inches(0.6),
             h, font=FONT_DISPLAY, size=24, color=INK)
    add_text(s, x, y0 + Inches(1.25), col_w, Inches(1.2),
             lead, font=FONT_BODY, size=13, color=INK, bold=True)
    add_text(s, x, y0 + Inches(2.25), col_w, Inches(1.8),
             body, font=FONT_BODY, size=11, color=MUTED)

footer(s, "APPROACH  ·  03")
page_number(s, 4, TOTAL)

# ---------- Slide 5 — Stats / quote -----------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PAPER)
add_text(s, Inches(0.8), Inches(0.7), Inches(4), Inches(0.4),
         "RESULTS", font=FONT_BODY, size=10, color=ACCENT,
         letter_spacing=600, bold=True)
add_line(s, Inches(0.8), Inches(1.15), Inches(0.5), Emu(15240), INK)
add_text(s, Inches(0.8), Inches(1.5), Inches(10), Inches(1.0),
         "Early signal from pilot customers",
         font=FONT_DISPLAY, size=34, color=INK)

# big stats row
stats = [("94%", "reduction in reconciliation time"),
         ("3.2×", "faster month-end close"),
         ("<40ms", "p95 query latency")]
x = Inches(0.8); y = Inches(3.1); w = Inches(3.9)
for i, (n, l) in enumerate(stats):
    xi = x + i * (w + Inches(0.15))
    add_text(s, xi, y, w, Inches(1.4),
             n, font=FONT_DISPLAY, size=64, color=INK)
    add_line(s, xi, y + Inches(1.55), Inches(0.4), Emu(15240), ACCENT_2)
    add_text(s, xi, y + Inches(1.8), w, Inches(0.8),
             l, font=FONT_BODY, size=12, color=MUTED)

# quote
add_rect(s, Inches(0.8), Inches(5.4), Emu(22860), Inches(1.4), ACCENT)
add_text(s, Inches(1.1), Inches(5.4), Inches(11), Inches(0.7),
         "\u201CIt replaced four tools and a spreadsheet. That never happens.\u201D",
         font=FONT_DISPLAY, size=18, color=INK)
add_text(s, Inches(1.1), Inches(6.15), Inches(11), Inches(0.4),
         "— Head of Finance Ops, Series C fintech",
         font=FONT_BODY, size=11, color=MUTED, letter_spacing=200)

footer(s, "RESULTS  ·  04")
page_number(s, 5, TOTAL)

# ---------- Slide 6 — Closing ------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
add_rect(s, Inches(0.8), Inches(1.0), Inches(0.6), Emu(22860), ACCENT)
add_text(s, Inches(0.8), Inches(1.2), Inches(6), Inches(0.4),
         "NEXT STEPS", font=FONT_BODY, size=10, color=ACCENT,
         letter_spacing=600, bold=True)
add_text(s, Inches(0.8), Inches(2.4), Inches(11), Inches(2.5),
         "Let's build\nsomething worth\nkeeping.",
         font=FONT_DISPLAY, size=60, color=PAPER)

add_text(s, Inches(0.8), Inches(5.9), Inches(4), Inches(0.35),
         "CONTACT", font=FONT_BODY, size=9, color=ACCENT,
         letter_spacing=500, bold=True)
add_text(s, Inches(0.8), Inches(6.3), Inches(5), Inches(0.4),
         "alanda.ambrose@aibos.co.jp",
         font=FONT_BODY, size=14, color=PAPER)

add_text(s, Inches(6.5), Inches(5.9), Inches(4), Inches(0.35),
         "NEXT", font=FONT_BODY, size=9, color=ACCENT,
         letter_spacing=500, bold=True)
add_text(s, Inches(6.5), Inches(6.3), Inches(6), Inches(0.4),
         "30-min walkthrough · Week of Apr 27",
         font=FONT_BODY, size=14, color=PAPER)

page_number(s, 6, TOTAL, color=HAIRLINE)

out = r"C:\Users\HP\Documents\PROJECTS\slide-desk\sample_premium.pptx"
prs.save(out)
print(f"Saved: {out}")
