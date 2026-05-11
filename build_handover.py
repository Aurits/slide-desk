"""
The Great Handover / 大承継
A bilingual executive deck on Japan's SMB succession opportunity.
Designed for Japanese investor audiences (JA primary, EN parallel).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Palette (inspired by traditional Japanese pigments) ---------------------
SUMI     = RGBColor(0x14, 0x14, 0x14)  # 墨 — ink black
KINARI   = RGBColor(0xF3, 0xEC, 0xDD)  # 生成 — unbleached paper
AI       = RGBColor(0x1D, 0x3B, 0x5C)  # 藍 — deep indigo
SHU      = RGBColor(0xB0, 0x3A, 0x2E)  # 朱 — vermillion (hanko red)
GIN      = RGBColor(0x9C, 0xA1, 0xA8)  # 銀 — silver-grey
CHA      = RGBColor(0x6B, 0x4F, 0x3A)  # 茶 — tea brown
HAIRLINE = RGBColor(0xC9, 0xBF, 0xA9)  # subtle divider on paper
FAINT    = RGBColor(0x3A, 0x3A, 0x3A)  # near-black secondary on paper

# --- Typography --------------------------------------------------------------
LATIN_DISPLAY = "Georgia"
LATIN_BODY    = "Segoe UI"
JA_DISPLAY    = "Yu Mincho"      # 游明朝 — serif, authoritative
JA_BODY       = "Yu Gothic UI"   # 游ゴシック — clean sans

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# --- XML helpers -------------------------------------------------------------
def set_ea_font(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)

def set_letter_spacing(run, spc):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(spc))

def set_vertical_ea(text_frame):
    bodyPr = text_frame._txBody.find(qn('a:bodyPr'))
    bodyPr.set('vert', 'eaVert')

# --- Primitives --------------------------------------------------------------
def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(s, color)
    return s

def hairline(slide, x, y, w, color=HAIRLINE, thickness=Emu(6350)):
    return rect(slide, x, y, w, thickness, color)

def vrule(slide, x, y, h, color=SHU, thickness=Emu(22860)):
    return rect(slide, x, y, thickness, h, color)

def text(slide, x, y, w, h, content, *,
         latin=LATIN_BODY, ja=JA_BODY, size=14, bold=False,
         color=SUMI, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content.split("\n") if isinstance(content, str) else content
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = latin
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        set_ea_font(r, ja)
        if spacing is not None:
            set_letter_spacing(r, spacing)
    return tb

def kicker(slide, x, y, w, jp, en, color=SHU):
    """Small bilingual label: JA · EN in caps, spaced."""
    text(slide, x, y, w, Inches(0.3), jp,
         latin=LATIN_BODY, ja=JA_BODY, size=10, bold=True,
         color=color, spacing=400)
    text(slide, x, y + Inches(0.32), w, Inches(0.3), en,
         latin=LATIN_BODY, ja=JA_BODY, size=9, bold=True,
         color=color, spacing=600)

def hanko(slide, x, y, size=Inches(0.9), char="信"):
    """Vermillion seal stamp — traditional Japanese signature mark."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = SHU
    s.line.color.rgb = SHU
    text(slide, x, y, size, size, char,
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=36, bold=True,
         color=KINARI, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def page_mark(slide, n, total, color=GIN):
    text(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.3),
         f"{n:02d} / {total:02d}",
         size=9, color=color, align=PP_ALIGN.RIGHT, spacing=400)

def footer_mark(slide, jp, en, color=GIN):
    text(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
         f"{jp}   ·   {en}", size=9, color=color, spacing=400, bold=True)

# --- Build -------------------------------------------------------------------
prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
blank = prs.slide_layouts[6]
TOTAL = 7

# ============================================================================
# SLIDE 1 — Cover
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
# left margin vertical indigo band (editorial anchor)
rect(s, 0, 0, Inches(0.45), SLIDE_H, AI)

# small kicker top
text(s, Inches(1.0), Inches(0.8), Inches(6), Inches(0.3),
     "AIBOS CAPITAL   ·   INVESTOR MEMO",
     size=10, color=AI, spacing=600, bold=True)
hairline(s, Inches(1.0), Inches(1.15), Inches(0.7), color=SHU, thickness=Emu(19050))

# Massive Japanese title
text(s, Inches(1.0), Inches(1.8), Inches(11), Inches(2.4),
     "大 承 継",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=108, bold=True,
     color=SUMI, spacing=400)

# English translation underneath
text(s, Inches(1.0), Inches(4.2), Inches(11), Inches(0.9),
     "The Great Handover",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=44, color=AI)

# Subtitle
text(s, Inches(1.0), Inches(5.2), Inches(10.5), Inches(0.5),
     "日本の中小企業承継危機を、世代の投資機会へ",
     size=14, color=FAINT, spacing=100)
text(s, Inches(1.0), Inches(5.55), Inches(11), Inches(0.5),
     "Reframing Japan's SMB succession crisis as a generational investment thesis.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, color=FAINT)

# Bottom line — date and confidentiality
hairline(s, Inches(1.0), Inches(6.8), Inches(11.3), color=HAIRLINE)
text(s, Inches(1.0), Inches(6.9), Inches(6), Inches(0.3),
     "2026 年 4 月   ·   April 2026",
     size=10, color=FAINT, spacing=300)
text(s, Inches(8.0), Inches(6.9), Inches(4), Inches(0.3),
     "CONFIDENTIAL   ·   機密",
     size=10, color=FAINT, align=PP_ALIGN.RIGHT, spacing=400, bold=True)

# Hanko seal
hanko(s, Inches(11.6), Inches(1.8), size=Inches(0.95), char="継")

# ============================================================================
# SLIDE 2 — The Moment (one striking statistic)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.7), Inches(0.7), Inches(6), "序章", "PROLOGUE")
hairline(s, Inches(0.7), Inches(1.45), Inches(11.9))

# Setup line
text(s, Inches(0.7), Inches(1.9), Inches(11), Inches(0.6),
     "Somewhere in Japan, every twenty minutes,",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, color=FAINT)
text(s, Inches(0.7), Inches(2.4), Inches(11), Inches(0.6),
     "a business closes — not because it failed, but because no one was ready to take over.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, color=FAINT)

# Huge statistic
text(s, Inches(0.7), Inches(3.4), Inches(12), Inches(2.2),
     "2,450,000",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=180, bold=True, color=SUMI)

# Stat caption
vrule(s, Inches(0.75), Inches(5.8), Inches(0.5), color=SHU, thickness=Emu(25400))
text(s, Inches(1.0), Inches(5.75), Inches(11), Inches(0.45),
     "後継者不在の日本企業 ／ 全中小企業の約三分の二",
     size=14, bold=True, color=SUMI, spacing=100)
text(s, Inches(1.0), Inches(6.2), Inches(11), Inches(0.5),
     "Japanese businesses without a designated successor — roughly two-thirds of all SMBs.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13, color=FAINT)

text(s, Inches(0.7), Inches(6.75), Inches(11), Inches(0.3),
     "Source: METI, Teikoku Databank (2023)",
     size=9, color=GIN, spacing=200)
footer_mark(s, "大承継", "THE GREAT HANDOVER")
page_mark(s, 2, TOTAL)

# ============================================================================
# SLIDE 3 — The Clock (demographic cliff, visualised)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.7), Inches(0.7), Inches(6), "第一章", "CHAPTER ONE")
hairline(s, Inches(0.7), Inches(1.45), Inches(11.9))

text(s, Inches(0.7), Inches(1.8), Inches(10), Inches(0.6),
     "時計は止まらない。",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, color=AI, bold=True)
text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.1),
     "The clock is not slowing down.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=40, color=SUMI, bold=True)

# Left column — narrative
text(s, Inches(0.7), Inches(4.0), Inches(5.7), Inches(0.5),
     "日本の町工場・老舗・地域企業。",
     size=13, bold=True, color=SUMI, spacing=100, line_spacing=1.3)
text(s, Inches(0.7), Inches(4.4), Inches(5.7), Inches(2.5),
     "Japan's workshops, century-old shops, and regional champions built the post-war economy. Their owners are now, on average, 62 years old — and 58% have no successor in place. Without a bridge, craft, IP, and local employment quietly evaporate.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, color=FAINT, line_spacing=1.4)

# Right column — age distribution visual
bar_x = Inches(7.2)
bar_w = Inches(5.4)
bar_top = Inches(4.0)
bar_h = Inches(0.42)
gap = Inches(0.18)

age_bands = [
    ("40代以下",   "Under 40", 0.08, CHA),
    ("50代",       "50s",      0.19, CHA),
    ("60代",       "60s",      0.34, AI),
    ("70代",       "70s",      0.26, SHU),
    ("80代以上",   "80+",      0.13, SUMI),
]
text(s, bar_x, Inches(3.55), bar_w, Inches(0.3),
     "SMB owner age distribution   ／   経営者年齢分布",
     size=10, bold=True, color=FAINT, spacing=200)

for i, (jp, en, pct, col) in enumerate(age_bands):
    y = bar_top + i * (bar_h + gap)
    text(s, bar_x, y, Inches(1.6), bar_h, f"{jp}  ·  {en}",
         size=10, color=SUMI, anchor=MSO_ANCHOR.MIDDLE, spacing=100)
    # track
    rect(s, bar_x + Inches(1.7), y + Inches(0.12),
         Inches(3.0), Inches(0.18), HAIRLINE)
    # fill
    rect(s, bar_x + Inches(1.7), y + Inches(0.12),
         Inches(3.0 * pct), Inches(0.18), col)
    text(s, bar_x + Inches(4.8), y, Inches(0.7), bar_h, f"{int(pct*100)}%",
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, bold=True,
         color=SUMI, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

text(s, bar_x, Inches(6.75), bar_w, Inches(0.3),
     "Source: SMBA White Paper (2024)",
     size=9, color=GIN, spacing=200)

footer_mark(s, "第一章 時計", "CHAPTER ONE · THE CLOCK")
page_mark(s, 3, TOTAL)

# ============================================================================
# SLIDE 4 — What's at Stake (dark, emotional)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, SUMI)
rect(s, 0, 0, Inches(0.45), SLIDE_H, SHU)

kicker(s, Inches(1.0), Inches(0.7), Inches(6),
       "第二章", "CHAPTER TWO", color=SHU)
hairline(s, Inches(1.0), Inches(1.45), Inches(11.3), color=RGBColor(0x2E, 0x2E, 0x2E))

text(s, Inches(1.0), Inches(1.9), Inches(11), Inches(1.0),
     "失われるのは、売上ではない。",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=26, color=GIN, bold=True)
text(s, Inches(1.0), Inches(2.45), Inches(11), Inches(1.1),
     "What disappears isn't revenue.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=44, color=KINARI, bold=True)

# Three stats in a row — scale of loss
stats = [
    ("¥22兆",     "22 trillion yen",     "GDP at risk by 2025\n2025年までに失われるGDP"),
    ("6.5M",       "six-and-a-half million", "Jobs in the balance\n雇用喪失の可能性"),
    ("127年",     "127 years",            "Avg. age of lost 老舗\n廃業老舗の平均創業年数"),
]
x0 = Inches(1.0); y0 = Inches(4.2); w = Inches(3.7); gap_x = Inches(0.2)
for i, (big, small, caption) in enumerate(stats):
    xi = x0 + i * (w + gap_x)
    text(s, xi, y0, w, Inches(1.2), big,
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=56, bold=True, color=KINARI)
    text(s, xi, y0 + Inches(1.3), w, Inches(0.35), small,
         size=10, color=SHU, spacing=500, bold=True)
    hairline(s, xi, y0 + Inches(1.7), Inches(0.5), color=SHU, thickness=Emu(12700))
    text(s, xi, y0 + Inches(1.85), w, Inches(1.0), caption,
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=11, color=GIN, line_spacing=1.4)

footer_mark(s, "第二章 代償", "CHAPTER TWO · THE COST", color=GIN)
page_mark(s, 4, TOTAL, color=GIN)

# ============================================================================
# SLIDE 5 — The Thesis (the turn in the story)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
# subtle indigo sidebar
rect(s, 0, 0, Inches(4.6), SLIDE_H, AI)

# Left panel (indigo)
text(s, Inches(0.6), Inches(0.7), Inches(4), Inches(0.3),
     "第三章",
     size=10, color=KINARI, spacing=600, bold=True)
text(s, Inches(0.6), Inches(1.0), Inches(4), Inches(0.3),
     "CHAPTER THREE",
     size=9, color=KINARI, spacing=600, bold=True)
hairline(s, Inches(0.6), Inches(1.5), Inches(0.6), color=SHU, thickness=Emu(19050))

text(s, Inches(0.6), Inches(2.0), Inches(4), Inches(1.8),
     "危機ではない。\n機会である。",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=32, color=KINARI, bold=True,
     line_spacing=1.2)
text(s, Inches(0.6), Inches(4.2), Inches(4), Inches(1.5),
     "Not a crisis.\nA portfolio.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=26, color=SHU, line_spacing=1.2)

# Right panel (paper)
text(s, Inches(5.2), Inches(0.9), Inches(7.5), Inches(0.5),
     "THE THESIS   ／   投資仮説",
     size=10, color=SHU, spacing=500, bold=True)
hairline(s, Inches(5.2), Inches(1.35), Inches(7), color=AI, thickness=Emu(12700))

text(s, Inches(5.2), Inches(1.7), Inches(7.5), Inches(1.8),
     "Japan's succession gap is the largest\nstructured private-market\nopportunity in Asia.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=26, color=SUMI, bold=True,
     line_spacing=1.2)

# Three pillars: Capital · Technology · Operators
pillars = [
    ("資本",    "Capital",     "Patient, long-duration vehicles that align with founder timelines.",
     "創業者の時間軸に寄り添う、長期保有型の資金。"),
    ("技術",    "Technology",  "A shared platform — ERP, finance, data — deployed across holdings.",
     "共通基盤による経営の高速化と可視化。"),
    ("人材",    "Operators",   "Bilingual operators who respect craft and can scale it.",
     "職人文化を理解し、拡張できるバイリンガル経営者。"),
]
yp = Inches(4.0)
for jp, en, body_en, body_jp in pillars:
    # JA big, EN under
    text(s, Inches(5.2), yp, Inches(1.6), Inches(0.5), jp,
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, bold=True, color=AI)
    text(s, Inches(5.2), yp + Inches(0.55), Inches(1.6), Inches(0.3), en.upper(),
         size=9, color=SHU, spacing=500, bold=True)
    text(s, Inches(7.0), yp, Inches(5.8), Inches(0.5), body_en,
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, bold=True, color=SUMI)
    text(s, Inches(7.0), yp + Inches(0.4), Inches(5.8), Inches(0.5), body_jp,
         size=11, color=FAINT, spacing=50)
    hairline(s, Inches(5.2), yp + Inches(0.85), Inches(7.5))
    yp += Inches(0.95)

footer_mark(s, "第三章 仮説", "CHAPTER THREE · THE THESIS")
page_mark(s, 5, TOTAL)

# ============================================================================
# SLIDE 6 — The Playbook: Shu-Ha-Ri
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.7), Inches(0.7), Inches(6), "第四章", "CHAPTER FOUR")
hairline(s, Inches(0.7), Inches(1.45), Inches(11.9))

text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.5),
     "守 · 破 · 離",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=34, bold=True, color=AI, spacing=800)
text(s, Inches(0.7), Inches(2.5), Inches(11), Inches(0.6),
     "The playbook, in the language of mastery.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, color=SUMI, bold=True)
text(s, Inches(0.7), Inches(3.05), Inches(11), Inches(0.5),
     "茶道・武道に由来する三段階 — 継承と革新の作法。",
     size=12, color=FAINT, spacing=100)

# Three stages
stages = [
    ("守", "SHU",  "Preserve",
     "Honour what works.",
     "First 12 months: retain the team, the recipe, the relationships. Map every invisible process the founder carries.",
     "最初の一年 — 人・商・縁を守り、暗黙知を可視化する。"),
    ("破", "HA",   "Adapt",
     "Modernise the engine.",
     "Year 2–3: deploy shared finance, data, and customer infrastructure. Remove the spreadsheet tax.",
     "二〜三年目 — 共通基盤を導入し、運営を再定義する。"),
    ("離", "RI",   "Scale",
     "Release from the mold.",
     "Year 4+: platform across acquisitions. The craft compounds; the overhead doesn't.",
     "四年目以降 — 複数企業の横展開。職人性は積み上がり、間接費は積み上がらない。"),
]
col_w = Inches(3.95)
x0 = Inches(0.7); y0 = Inches(3.8); gap_x = Inches(0.2)

for i, (kanji, romaji, english, head_en, body_en, body_jp) in enumerate(stages):
    xi = x0 + i * (col_w + gap_x)
    # Kanji card background
    rect(s, xi, y0, col_w, Inches(3.0), KINARI)
    # Big kanji
    text(s, xi, y0 + Inches(0.05), Inches(1.2), Inches(1.2), kanji,
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=72, bold=True,
         color=SHU if i == 1 else AI)
    # Romaji + English
    text(s, xi + Inches(1.3), y0 + Inches(0.15), Inches(2.6), Inches(0.35),
         romaji, size=10, color=SHU, spacing=600, bold=True)
    text(s, xi + Inches(1.3), y0 + Inches(0.5), Inches(2.6), Inches(0.45),
         english, latin=LATIN_DISPLAY, ja=JA_DISPLAY,
         size=20, bold=True, color=SUMI)
    # Headline
    text(s, xi, y0 + Inches(1.35), col_w, Inches(0.45),
         head_en, latin=LATIN_DISPLAY, ja=JA_DISPLAY,
         size=15, bold=True, color=AI)
    hairline(s, xi, y0 + Inches(1.85), Inches(0.5),
             color=SHU, thickness=Emu(12700))
    # Body EN
    text(s, xi, y0 + Inches(2.0), col_w, Inches(1.3),
         body_en, latin=LATIN_DISPLAY, ja=JA_DISPLAY,
         size=11, color=FAINT, line_spacing=1.4)
    # Body JP
    text(s, xi, y0 + Inches(3.1), col_w, Inches(0.6),
         body_jp, size=10, color=FAINT, line_spacing=1.4)

footer_mark(s, "第四章 守破離", "CHAPTER FOUR · SHU · HA · RI")
page_mark(s, 6, TOTAL)

# ============================================================================
# SLIDE 7 — Proof + Ask (closing)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, AI)
rect(s, 0, 0, Inches(0.45), SLIDE_H, SHU)

kicker(s, Inches(1.0), Inches(0.7), Inches(6),
       "終章", "EPILOGUE", color=KINARI)
hairline(s, Inches(1.0), Inches(1.45), Inches(11.3),
         color=RGBColor(0x2E, 0x4F, 0x7A))

text(s, Inches(1.0), Inches(1.85), Inches(11), Inches(1.0),
     "静かな革命に、",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=26, color=GIN, bold=True)
text(s, Inches(1.0), Inches(2.4), Inches(11), Inches(1.3),
     "Join the quiet revolution.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=48, color=KINARI, bold=True)

# Early proof stat (left)
text(s, Inches(1.0), Inches(4.2), Inches(5.5), Inches(0.35),
     "PILOT SIGNAL   ／   初期実績",
     size=10, color=SHU, spacing=500, bold=True)
hairline(s, Inches(1.0), Inches(4.6), Inches(1.0),
         color=SHU, thickness=Emu(19050))
text(s, Inches(1.0), Inches(4.8), Inches(5.5), Inches(1.2),
     "+38% EBITDA",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=54, bold=True, color=KINARI)
text(s, Inches(1.0), Inches(5.85), Inches(5.5), Inches(0.4),
     "across 3 portfolio companies, 18 months post-transition",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=11, color=GIN)
text(s, Inches(1.0), Inches(6.2), Inches(5.5), Inches(0.4),
     "3社平均・承継後18ヶ月",
     size=10, color=GIN, spacing=100)

# Ask (right)
text(s, Inches(7.2), Inches(4.2), Inches(5.5), Inches(0.35),
     "THE ASK   ／   ご提案",
     size=10, color=SHU, spacing=500, bold=True)
hairline(s, Inches(7.2), Inches(4.6), Inches(1.0),
         color=SHU, thickness=Emu(19050))
text(s, Inches(7.2), Inches(4.8), Inches(5.5), Inches(0.5),
     "Anchor commitment · ¥5B",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, bold=True, color=KINARI)
text(s, Inches(7.2), Inches(5.25), Inches(5.5), Inches(0.4),
     "アンカー出資 ／ 50億円",
     size=12, color=GIN, spacing=100)
text(s, Inches(7.2), Inches(5.8), Inches(5.5), Inches(0.5),
     "Advisory seat · Fund II close Q3 2026",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13, color=KINARI)
text(s, Inches(7.2), Inches(6.2), Inches(5.5), Inches(0.4),
     "アドバイザリー席・2026年第3四半期クローズ",
     size=10, color=GIN, spacing=100)

# Signature line
hairline(s, Inches(1.0), Inches(6.85), Inches(11.3),
         color=RGBColor(0x2E, 0x4F, 0x7A))
text(s, Inches(1.0), Inches(6.95), Inches(7), Inches(0.3),
     "alanda.ambrose @ aibos.co.jp",
     size=10, color=GIN, spacing=200)
text(s, Inches(10.5), Inches(6.95), Inches(2.3), Inches(0.3),
     f"{7:02d} / {TOTAL:02d}",
     size=9, color=GIN, align=PP_ALIGN.RIGHT, spacing=400)

# Hanko on closing
hanko(s, Inches(11.5), Inches(1.85), size=Inches(0.95), char="承")

# --- Save --------------------------------------------------------------------
out = r"C:\Users\HP\Documents\PROJECTS\slide-desk\great_handover.pptx"
prs.save(out)
print(f"Saved: {out}")
