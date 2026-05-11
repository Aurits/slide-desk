"""
A Note from Create Project — a quarterly update for the board.
Seventeen slides: honest, concrete, creative.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Palette · Hand & Rule ---------------------------------------------------
# Names kept for code continuity; semantic roles:
#   SUMI    → Soot (text)
#   KINARI  → Parchment (paper)
#   KINARI_D→ Parchment soft (raised surface)
#   AI      → Rule (Blueprint's domain colour)
#   AI_SOFT → Rule soft (lighter rule for dark-slide accents)
#   SHU     → Thread (Apparel's domain colour)
#   BRASS   → Brass (highlight / second accent)
#   GIN     → Ash (secondary neutral)
SUMI     = RGBColor(0x15, 0x12, 0x11)   # Soot
KINARI   = RGBColor(0xF4, 0xEE, 0xE3)   # Parchment
KINARI_D = RGBColor(0xEC, 0xE4, 0xD3)   # Parchment soft
AI       = RGBColor(0x25, 0x48, 0x72)   # Rule
AI_SOFT  = RGBColor(0x44, 0x65, 0x8C)   # Rule soft
SHU      = RGBColor(0xB1, 0x50, 0x30)   # Thread
BRASS    = RGBColor(0xC2, 0x90, 0x46)   # Brass
GIN      = RGBColor(0x89, 0x83, 0x78)   # Ash
HAIRLINE = RGBColor(0xD5, 0xCA, 0xB4)
FAINT    = RGBColor(0x3A, 0x35, 0x30)

# Fonts — Sitka (display) + Corbel (body) on Windows; Yu Mincho/Gothic for JP
LATIN_DISPLAY = "Sitka Banner"
LATIN_BODY    = "Corbel"
JA_DISPLAY    = "Yu Mincho"
JA_BODY       = "Yu Gothic UI"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# --- XML helpers -------------------------------------------------------------
def set_ea_font(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)

def set_letter_spacing(run, spc):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(spc))

def set_transition(slide, kind='fade'):
    sld = slide._element
    for t in sld.findall(qn('p:transition')):
        sld.remove(t)
    t = etree.SubElement(sld, qn('p:transition'))
    t.set('spd', 'med')
    etree.SubElement(t, qn(f'p:{kind}'))

def animate_on_click(slide, groups):
    if not groups:
        return
    group_ids = [[sh.shape_id for sh in g] for g in groups]
    sld = slide._element
    for t in sld.findall(qn('p:timing')):
        sld.remove(t)

    timing = etree.SubElement(sld, qn('p:timing'))
    tnLst = etree.SubElement(timing, qn('p:tnLst'))
    root_par = etree.SubElement(tnLst, qn('p:par'))
    root_cTn = etree.SubElement(root_par, qn('p:cTn'))
    root_cTn.set('id', '1'); root_cTn.set('dur', 'indefinite')
    root_cTn.set('restart', 'never'); root_cTn.set('nodeType', 'tmRoot')
    root_childTnLst = etree.SubElement(root_cTn, qn('p:childTnLst'))

    seq = etree.SubElement(root_childTnLst, qn('p:seq'))
    seq.set('concurrent', '1'); seq.set('nextAc', 'seek')
    seq_cTn = etree.SubElement(seq, qn('p:cTn'))
    seq_cTn.set('id', '2'); seq_cTn.set('dur', 'indefinite')
    seq_cTn.set('nodeType', 'mainSeq')
    seq_childTnLst = etree.SubElement(seq_cTn, qn('p:childTnLst'))

    counter = [3]
    def nid():
        v = counter[0]; counter[0] += 1; return str(v)

    for group in group_ids:
        click_par = etree.SubElement(seq_childTnLst, qn('p:par'))
        click_cTn = etree.SubElement(click_par, qn('p:cTn'))
        click_cTn.set('id', nid()); click_cTn.set('fill', 'hold')
        stCondLst = etree.SubElement(click_cTn, qn('p:stCondLst'))
        etree.SubElement(stCondLst, qn('p:cond')).set('delay', 'indefinite')
        click_childTnLst = etree.SubElement(click_cTn, qn('p:childTnLst'))

        par2 = etree.SubElement(click_childTnLst, qn('p:par'))
        cTn2 = etree.SubElement(par2, qn('p:cTn'))
        cTn2.set('id', nid()); cTn2.set('fill', 'hold')
        stCondLst2 = etree.SubElement(cTn2, qn('p:stCondLst'))
        etree.SubElement(stCondLst2, qn('p:cond')).set('delay', '0')
        childTnLst2 = etree.SubElement(cTn2, qn('p:childTnLst'))

        for i, sid in enumerate(group):
            par3 = etree.SubElement(childTnLst2, qn('p:par'))
            cTn3 = etree.SubElement(par3, qn('p:cTn'))
            cTn3.set('id', nid()); cTn3.set('presetID', '10')
            cTn3.set('presetClass', 'entr'); cTn3.set('presetSubtype', '0')
            cTn3.set('fill', 'hold'); cTn3.set('grpId', '0')
            cTn3.set('nodeType', 'clickEffect' if i == 0 else 'withEffect')
            stCondLst3 = etree.SubElement(cTn3, qn('p:stCondLst'))
            etree.SubElement(stCondLst3, qn('p:cond')).set('delay', '0')
            childTnLst3 = etree.SubElement(cTn3, qn('p:childTnLst'))

            set_el = etree.SubElement(childTnLst3, qn('p:set'))
            cBhvr = etree.SubElement(set_el, qn('p:cBhvr'))
            cTn_b = etree.SubElement(cBhvr, qn('p:cTn'))
            cTn_b.set('id', nid()); cTn_b.set('dur', '1'); cTn_b.set('fill', 'hold')
            stCondLst_b = etree.SubElement(cTn_b, qn('p:stCondLst'))
            etree.SubElement(stCondLst_b, qn('p:cond')).set('delay', '0')
            tgtEl = etree.SubElement(cBhvr, qn('p:tgtEl'))
            etree.SubElement(tgtEl, qn('p:spTgt')).set('spid', str(sid))
            attrNameLst = etree.SubElement(cBhvr, qn('p:attrNameLst'))
            etree.SubElement(attrNameLst, qn('p:attrName')).text = 'style.visibility'
            to = etree.SubElement(set_el, qn('p:to'))
            etree.SubElement(to, qn('p:strVal')).set('val', 'visible')

            animEffect = etree.SubElement(childTnLst3, qn('p:animEffect'))
            animEffect.set('transition', 'in'); animEffect.set('filter', 'fade')
            cBhvr_f = etree.SubElement(animEffect, qn('p:cBhvr'))
            cTn_f = etree.SubElement(cBhvr_f, qn('p:cTn'))
            cTn_f.set('id', nid()); cTn_f.set('dur', '700')
            tgtEl_f = etree.SubElement(cBhvr_f, qn('p:tgtEl'))
            etree.SubElement(tgtEl_f, qn('p:spTgt')).set('spid', str(sid))

    prevCondLst = etree.SubElement(seq, qn('p:prevCondLst'))
    pcond = etree.SubElement(prevCondLst, qn('p:cond'))
    pcond.set('evt', 'onPrev'); pcond.set('delay', '0')
    etree.SubElement(etree.SubElement(pcond, qn('p:tgtEl')), qn('p:sldTgt'))
    nextCondLst = etree.SubElement(seq, qn('p:nextCondLst'))
    ncond = etree.SubElement(nextCondLst, qn('p:cond'))
    ncond.set('evt', 'onNext'); ncond.set('delay', '0')
    etree.SubElement(etree.SubElement(ncond, qn('p:tgtEl')), qn('p:sldTgt'))

# --- Drawing primitives ------------------------------------------------------
def fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h); fill(s, color); return s

def hairline(slide, x, y, w, color=HAIRLINE, thickness=Emu(6350)):
    return rect(slide, x, y, w, thickness, color)

def oval(slide, x, y, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size); fill(s, color); return s

def text(slide, x, y, w, h, content, *,
         latin=LATIN_BODY, ja=JA_BODY, size=14, bold=False, italic=False,
         color=SUMI, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content.split("\n") if isinstance(content, str) else content
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = latin; r.font.size = Pt(size)
        r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
        set_ea_font(r, ja)
        if spacing is not None: set_letter_spacing(r, spacing)
    return tb

def kicker(slide, x, y, w, jp, en, color=SHU):
    a = text(slide, x, y, w, Inches(0.3), jp, size=10, bold=True, color=color, spacing=400)
    b = text(slide, x, y + Inches(0.32), w, Inches(0.3), en, size=9, bold=True, color=color, spacing=600)
    return a, b

def hanko(slide, x, y, size=Inches(0.9), char="創"):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = SHU; s.line.color.rgb = SHU
    t = text(slide, x, y, size, size, char,
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=34, bold=True,
             color=KINARI, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s, t

def page_mark(slide, n, total, color=GIN):
    return text(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.3),
                f"{n:02d} / {total:02d}",
                size=9, color=color, align=PP_ALIGN.RIGHT, spacing=400)

def footer_mark(slide, jp, en, color=GIN):
    return text(slide, Inches(0.7), Inches(7.05), Inches(9), Inches(0.3),
                f"{jp}   ·   {en}",
                size=9, color=color, spacing=400, bold=True)

# --- Build -------------------------------------------------------------------
prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
blank = prs.slide_layouts[6]
TOTAL = 17

# ============================================================================
# 1 — Cover
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
rect(s, 0, 0, Inches(0.4), SLIDE_H, AI)

text(s, Inches(0.9), Inches(0.8), Inches(6), Inches(0.3),
     "AIBOS   ·   CREATE PROJECT",
     size=10, color=AI, spacing=600, bold=True)
hairline(s, Inches(0.9), Inches(1.15), Inches(0.7), color=SHU, thickness=Emu(19050))

text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.3),
     "A Note from",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=58, color=AI, italic=True)
text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.8),
     "Create Project",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=86, bold=True, color=SUMI)

hairline(s, Inches(0.9), Inches(4.9), Inches(0.7), color=SHU, thickness=Emu(19050))
text(s, Inches(0.9), Inches(5.05), Inches(11.5), Inches(0.5),
     "Two products, a team of six, and the quarter behind them.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=18, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.5),
     "二つのプロダクト、六人のチーム、そしてこの四半期の歩み。",
     size=12, color=FAINT, spacing=100)

text(s, Inches(0.9), Inches(6.3), Inches(6), Inches(0.3),
     "クリエイト・プロジェクト便り",
     size=11, color=AI, italic=True, spacing=200)

hairline(s, Inches(0.9), Inches(6.8), Inches(11.4))
text(s, Inches(0.9), Inches(6.9), Inches(6), Inches(0.3),
     "Q2 · 2026年春   ·   Spring 2026",
     size=10, color=FAINT, spacing=300)
text(s, Inches(7.5), Inches(6.9), Inches(4.8), Inches(0.3),
     "FOR THE BOARD   ·   取締役会用",
     size=10, color=FAINT, align=PP_ALIGN.RIGHT, spacing=500, bold=True)

hanko(s, Inches(11.6), Inches(1.7), size=Inches(0.95), char="創")
set_transition(s)

# ============================================================================
# 2 — Opening note (three-line reveal)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6), "ご挨拶", "A NOTE")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

line1 = text(s, Inches(1.5), Inches(2.6), Inches(11), Inches(1.0),
             "A short note from Create Project —",
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=36, color=SUMI, italic=True)
line2 = text(s, Inches(1.5), Inches(3.4), Inches(11), Inches(1.0),
             "on what the team has been building this quarter,",
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=36, color=SUMI, italic=True)
line3 = text(s, Inches(1.5), Inches(4.4), Inches(11), Inches(1.0),
             "and where it is going next.",
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=36, bold=True, color=AI)

hairline(s, Inches(1.5), Inches(5.8), Inches(0.8), color=SHU, thickness=Emu(19050))
jp_line = text(s, Inches(1.5), Inches(6.05), Inches(10), Inches(0.5),
               "今期に作ってきたもの、そしてこれからの行き先について、一筆。",
               size=14, color=FAINT, italic=True, spacing=100)

footer_mark(s, "ご挨拶", "A NOTE")
page_mark(s, 2, TOTAL)
animate_on_click(s, [[line1], [line2], [line3, jp_line]])
set_transition(s)

# ============================================================================
# 3 — The team
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6), "仲間", "THE TEAM")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.9),
     "Six makers.  Two products.  One team.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=38, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.7), Inches(11), Inches(0.45),
     "Every strength is known — every role assigned accordingly.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15, italic=True, color=SUMI)
text(s, Inches(0.9), Inches(3.1), Inches(11), Inches(0.4),
     "一人ひとりの強みを把握し、役割を的確に振り分けています。",
     size=12, color=FAINT, italic=True, spacing=100)

row1 = [
    ("Ambrose",  "Systems & Product",  "統率"),
    ("Aita",     "Workflow Engine",    "設計"),
    ("Muzuva",   "Backend & Pipeline", "基盤"),
]
row2 = [
    ("Lauben",   "Infrastructure",     "土台"),
    ("Crispus",  "Frontend",           "表層"),
    ("Otai",     "Backend Support",    "補佐"),
]

def plot_row(people, y_dot, x_start, x_step, dot_color=SHU):
    for i, (name, role, jp) in enumerate(people):
        cx = Inches(x_start + i * x_step)
        oval(s, cx, y_dot, Inches(0.26), dot_color)
        text(s, cx - Inches(1.4), y_dot - Inches(0.5), Inches(3.0), Inches(0.3),
             jp, size=11, color=AI, align=PP_ALIGN.CENTER, spacing=300, bold=True)
        text(s, cx - Inches(1.4), y_dot + Inches(0.32), Inches(3.0), Inches(0.3),
             name, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15, bold=True,
             color=SUMI, align=PP_ALIGN.CENTER)
        text(s, cx - Inches(1.6), y_dot + Inches(0.7), Inches(3.4), Inches(0.3),
             role, size=11, color=FAINT, align=PP_ALIGN.CENTER, italic=True)

plot_row(row1, Inches(4.05), 2.5, 4.17)   # 2.50, 6.67, 10.84
plot_row(row2, Inches(5.65), 2.5, 4.17)

hairline(s, Inches(0.9), Inches(6.6), Inches(11.5))
text(s, Inches(0.9), Inches(6.75), Inches(11), Inches(0.4),
     "Engineering, product, infrastructure — all carried by the same six people.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13, italic=True, color=FAINT)

footer_mark(s, "仲間", "THE TEAM")
page_mark(s, 3, TOTAL)
set_transition(s)

# ============================================================================
# 4 — How we build (three principles, specific to Create Project)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6), "流儀", "HOW WE BUILD")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.8),
     "Three principles every feature is held to.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=26, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.45), Inches(11), Inches(0.4),
     "すべての機能が通る、三つの原則。",
     size=12, color=FAINT, spacing=100)

principles = [
    ("Traceability", "追跡", "Every feature",
     "maps back to a customer's exact words — no code without a request.",
     "顧客の言葉に紐づかない機能は、書かない。"),
    ("Ownership", "責任", "The engineer who ships it",
     "is the engineer who stands behind it when it runs in production.",
     "作る者が、運用を支える者になる。"),
    ("Design first", "設計", "We prove the design",
     "on paper before we write a line of code — rework starts before the build.",
     "コードを書く前に、設計を確かめる。"),
]
col_w = Inches(3.8); gap = Inches(0.3); x0 = Inches(0.9); y0 = Inches(3.25)
col_groups = []
for i, (en, jp, head, body, jp_body) in enumerate(principles):
    x = x0 + i * (col_w + gap)
    a = text(s, x, y0, col_w, Inches(0.6), en,
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=28, bold=True, color=AI)
    b = text(s, x, y0 + Inches(0.65), col_w, Inches(0.5), jp,
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, color=SHU, spacing=400)
    c = hairline(s, x, y0 + Inches(1.3), Inches(0.5), color=SHU, thickness=Emu(12700))
    d = text(s, x, y0 + Inches(1.45), col_w, Inches(0.5), head,
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15, bold=True,
             color=SUMI, italic=True)
    e = text(s, x, y0 + Inches(1.85), col_w, Inches(1.5), body,
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13, italic=True,
             color=FAINT, line_spacing=1.5)
    f = text(s, x, y0 + Inches(3.1), col_w, Inches(0.5), jp_body,
             size=10, color=FAINT, spacing=100, line_spacing=1.4)
    col_groups.append([a, b, c, d, e, f])

footer_mark(s, "流儀", "HOW WE BUILD")
page_mark(s, 4, TOTAL)
animate_on_click(s, col_groups)
set_transition(s)

# ============================================================================
# 5 — Apparel · the scene
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
rect(s, Inches(0.9), Inches(0.7), Emu(22860), Inches(6.3), SHU)

text(s, Inches(1.3), Inches(0.7), Inches(5), Inches(0.35),
     "其ノ壱   ·   APPAREL",
     size=10, color=SHU, spacing=600, bold=True)
text(s, Inches(1.3), Inches(1.15), Inches(11), Inches(1.1),
     "Apparel",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=64, bold=True, color=SUMI)
text(s, Inches(1.3), Inches(2.15), Inches(11), Inches(0.5),
     "A product photo becomes a full fashion campaign.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15, italic=True, color=FAINT)
text(s, Inches(1.3), Inches(2.55), Inches(11), Inches(0.5),
     "一枚の商品写真から、プロ仕様のキャンペーンへ。",
     size=12, color=FAINT, spacing=100)

hairline(s, Inches(1.3), Inches(3.4), Inches(0.7),
         color=SHU, thickness=Emu(19050))

v1 = text(s, Inches(1.3), Inches(3.65), Inches(10.5), Inches(0.6),
          "A brand uploads one photo of a garment.",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, italic=True, color=SUMI)
v2 = text(s, Inches(1.3), Inches(4.35), Inches(10.5), Inches(0.6),
          "The team picks the model, the background, the mood.",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, italic=True, color=SUMI)
v3 = text(s, Inches(1.3), Inches(5.05), Inches(10.5), Inches(0.6),
          "In the time it takes to pour a coffee, a full campaign arrives —",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, italic=True, color=SUMI)
v4 = text(s, Inches(1.3), Inches(5.75), Inches(10.5), Inches(0.6),
          "every look, every angle, ready to publish.",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, bold=True, color=AI)

footer_mark(s, "其ノ壱 · APPAREL", "PRODUCT ONE")
page_mark(s, 5, TOTAL)
animate_on_click(s, [[v1], [v2], [v3], [v4]])
set_transition(s)

# ============================================================================
# 6 — Apparel · in action (image)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6), "実例", "APPAREL · IN ACTION")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "Apparel, in action.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=36, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.4),
     "The landing showcase, and the tool the team actually uses.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.3),
     "紹介ページと、実際のツール画面。",
     size=11, color=FAINT, italic=True, spacing=100)

# Stacked image cards — back (landing) behind, front (tool) on top
SHADOW_BACK  = RGBColor(0xBE, 0xB3, 0x97)
SHADOW_FRONT = RGBColor(0x8F, 0x83, 0x68)

# Actual image aspects (h/w): apparel.png 0.428, apparel1.png 0.432
# After crop_top=0.13 on apparel.png, visible aspect = 0.428 × 0.87 ≈ 0.372

# BACK — apparel.png (landing showcase) — upper-right, smaller
bw, bx, by = 5.2, 7.1, 2.95
bh = bw * 0.372
rect(s, Inches(bx + 0.07), Inches(by + 0.08),
     Inches(bw), Inches(bh), SHADOW_BACK)
pic_back = s.shapes.add_picture(
    r"C:\Users\HP\Documents\PROJECTS\slide-desk\apparel.png",
    Inches(bx), Inches(by), width=Inches(bw), height=Inches(bh))
pic_back.crop_top = 0.13

# FRONT — apparel1.png (the tool interface) — lower-left, hero
fw, fx, fy = 7.0, 0.9, 3.85
fh = fw * 0.432
rect(s, Inches(fx + 0.09), Inches(fy + 0.11),
     Inches(fw), Inches(fh), SHADOW_FRONT)
pic_front = s.shapes.add_picture(
    r"C:\Users\HP\Documents\PROJECTS\slide-desk\apparel1.png",
    Inches(fx), Inches(fy), width=Inches(fw), height=Inches(fh))

footer_mark(s, "実例", "APPAREL · IN ACTION")
page_mark(s, 6, TOTAL)
set_transition(s)

# ============================================================================
# 7 — Apparel · now and next (operational)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "進捗 · 其ノ壱", "APPAREL · NOW & NEXT")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "Apparel — now and next.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=34, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
     "What's running, what's coming, who's doing it.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.3),
     "今、進めていること。次に取り掛かること。担うのは誰か。",
     size=11, color=FAINT, italic=True, spacing=100)

# Two columns: NOW · NEXT (team is on the next slide)
col_w = Inches(5.4); col_y = Inches(3.45)

def draw_column(slide, x, header, header_color, items):
    shapes = []
    shapes.append(text(slide, x, col_y, col_w, Inches(0.32),
         header, size=10, bold=True, color=header_color, spacing=500))
    shapes.append(hairline(slide, x, col_y + Inches(0.42), Inches(0.55),
             color=header_color, thickness=Emu(19050)))
    iy = col_y + Inches(0.75)
    for en_line, jp_line in items:
        shapes.append(oval(slide, x, iy + Inches(0.1), Inches(0.09), header_color))
        shapes.append(text(slide, x + Inches(0.25), iy, Inches(5.0), Inches(0.4),
             en_line, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14,
             color=SUMI, bold=True, line_spacing=1.3))
        shapes.append(text(slide, x + Inches(0.25), iy + Inches(0.4), Inches(5.0),
             Inches(0.35), jp_line,
             size=11, color=FAINT, italic=True, spacing=50, line_spacing=1.3))
        iy += Inches(0.85)
    return shapes

x_now  = Inches(0.9)
x_next = Inches(0.9 + 5.4 + 0.7)   # 7.0

now_shapes = draw_column(s, x_now, "NOW · 進行中", GIN, [
    ("Production monitoring + GCP cost oversight",
     "本番運用とGCP費用の監視"),
    ("Multi-angle background consistency — R&D",
     "多角度背景の一貫性 — 検証中"),
    ("Client-driven UX adjustments",
     "顧客要望によるUX調整"),
    ("ComfyUI pipeline maintenance",
     "ComfyUIパイプラインの保守"),
])

next_shapes = draw_column(s, x_next, "NEXT · 60 DAYS · 次の60日", SHU, [
    ("Sharper background engine — new perspectives, richer scenes",
     "背景エンジンの精度と視点を一段引き上げる"),
    ("Onboard new paying customers",
     "新規有償顧客のオンボーディング"),
    ("Document ComfyUI pipeline; assign a backup engineer",
     "ComfyUIパイプラインの文書化と予備担当の任命"),
])

footer_mark(s, "進捗 · APPAREL", "NOW & NEXT")
page_mark(s, 7, TOTAL)
animate_on_click(s, [now_shapes, next_shapes])
set_transition(s)

# ============================================================================
# 8 — Apparel · the team (Muzuva + Crispus)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "体制 · 其ノ壱", "APPAREL · THE TEAM")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "The Apparel team.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=34, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
     "A small dedicated pair, focused on the work.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.3),
     "アパレルを支える、二人組。",
     size=11, color=FAINT, italic=True, spacing=100)

# Two large cards centered
card_w = Inches(5.0); card_h = Inches(2.7); gap_x = Inches(0.4)
total = 5.0 * 2 + 0.4
left_margin = (13.333 - total) / 2
cx1 = Inches(left_margin)
cx2 = cx1 + card_w + gap_x
cy_card = Inches(3.5)

def apparel_card(x, kanji, name, role, hours, note_en, note_jp):
    rect(s, x, cy_card, card_w, card_h, KINARI_D)
    rect(s, x, cy_card, card_w, Emu(28575), SHU)  # top bar in Thread
    # Kanji top-right
    text(s, x + card_w - Inches(1.2), cy_card + Inches(0.25),
         Inches(1.0), Inches(0.7),
         kanji, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=46,
         color=SHU, bold=True, align=PP_ALIGN.RIGHT)
    # Name
    text(s, x + Inches(0.4), cy_card + Inches(0.3),
         card_w - Inches(1.4), Inches(0.7),
         name, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=34,
         bold=True, color=SUMI)
    # Role
    text(s, x + Inches(0.4), cy_card + Inches(0.95),
         card_w - Inches(0.8), Inches(0.4),
         role, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14,
         italic=True, color=AI)
    # Hours
    text(s, x + Inches(0.4), cy_card + Inches(1.3),
         card_w - Inches(0.8), Inches(0.5),
         hours, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22,
         bold=True, color=BRASS)
    # Hairline
    hairline(s, x + Inches(0.4), cy_card + Inches(1.85),
             Inches(0.5), color=SHU, thickness=Emu(15240))
    # Note
    text(s, x + Inches(0.4), cy_card + Inches(2.0),
         card_w - Inches(0.8), Inches(0.35),
         note_en, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=11,
         color=SUMI, italic=True, line_spacing=1.3)
    text(s, x + Inches(0.4), cy_card + Inches(2.32),
         card_w - Inches(0.8), Inches(0.35),
         note_jp, size=10, color=FAINT, italic=True, spacing=50)

apparel_card(cx1, "基盤", "Muzuva", "Backend & ComfyUI Lead", "≈ 16 h / week",
             "Splits time with Blueprint backend (16 h there).",
             "ブループリント側にも16時間を担当。")
apparel_card(cx2, "表層", "Crispus", "Frontend", "≈ 8 h / week",
             "Splits time with Blueprint frontend (32 h there).",
             "ブループリント側に32時間、こちらは8時間。")

# Bottom summary
hairline(s, Inches(0.9), Inches(6.5), Inches(11.5))
text(s, Inches(0.9), Inches(6.62), Inches(11.5), Inches(0.35),
     "2 contributors  ·  24 hrs / week on Apparel  ·  Both also support Blueprint — flexible cross-product team.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, bold=True, color=AI)
text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.3),
     "アパレル側は2名・24時間/週。両者ともブループリントにも従事し、柔軟に横断。",
     size=10, color=FAINT, italic=True, spacing=100)

footer_mark(s, "体制 · APPAREL", "THE TEAM")
page_mark(s, 8, TOTAL)
set_transition(s)

# ============================================================================
# 8 — Blueprint · the scene
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, SUMI)
rect(s, Inches(0.9), Inches(0.7), Emu(22860), Inches(6.3), AI_SOFT)

text(s, Inches(1.3), Inches(0.7), Inches(5), Inches(0.35),
     "其ノ弐   ·   BLUEPRINT",
     size=10, color=AI_SOFT, spacing=600, bold=True)
text(s, Inches(1.3), Inches(1.15), Inches(11), Inches(1.1),
     "Blueprint",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=64, bold=True, color=KINARI)
text(s, Inches(1.3), Inches(2.15), Inches(11), Inches(0.5),
     "Upload a blueprint.  Get instant analysis, parts list, and inspection.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15, italic=True, color=HAIRLINE)
text(s, Inches(1.3), Inches(2.55), Inches(11), Inches(0.5),
     "図面を投げ込むだけで、解析・部品表・検査が返ってくる。",
     size=12, color=HAIRLINE, spacing=100)

hairline(s, Inches(1.3), Inches(3.4), Inches(0.7),
         color=AI_SOFT, thickness=Emu(19050))

v1 = text(s, Inches(1.3), Inches(3.65), Inches(10.5), Inches(0.6),
          "An engineer drops a technical drawing into Blueprint.",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, italic=True, color=KINARI)
v2 = text(s, Inches(1.3), Inches(4.35), Inches(10.5), Inches(0.6),
          "Minutes later, three answers come back —",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, italic=True, color=KINARI)
v3 = text(s, Inches(1.3), Inches(5.05), Inches(10.5), Inches(0.6),
          "the parts list, already tabled.  The errors, circled with confidence.",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, italic=True, color=KINARI)
v4 = text(s, Inches(1.3), Inches(5.75), Inches(10.5), Inches(0.6),
          "The text, searchable as any document.",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, bold=True, color=SHU)

footer_mark(s, "其ノ弐 · BLUEPRINT", "PRODUCT TWO", color=HAIRLINE)
page_mark(s, 9, TOTAL, color=HAIRLINE)
animate_on_click(s, [[v1], [v2], [v3], [v4]])
set_transition(s)

# ============================================================================
# 9 — Blueprint · at first glance (image)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "動作中", "BLUEPRINT · AT WORK", color=AI)
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "Blueprint, at work.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=36, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.4),
     "From the entry page to the parts list it extracts from a drawing.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.3),
     "入口から、図面より抽出される構造化された部品表まで。",
     size=11, color=FAINT, italic=True, spacing=100)

# Stacked image cards — back (entry page) behind, front (BOM extraction) on top
SHADOW_BACK  = RGBColor(0xBE, 0xB3, 0x97)
SHADOW_FRONT = RGBColor(0x8F, 0x83, 0x68)

# Actual image aspects (h/w): bp.png 0.428, bp1.png 0.428

# BACK — bp.png (entry page) — upper-right, smaller
bw, bx, by = 5.0, 7.3, 2.95
bh = bw * 0.428
rect(s, Inches(bx + 0.07), Inches(by + 0.08),
     Inches(bw), Inches(bh), SHADOW_BACK)
pic_back = s.shapes.add_picture(
    r"C:\Users\HP\Documents\PROJECTS\slide-desk\bp.png",
    Inches(bx), Inches(by), width=Inches(bw), height=Inches(bh))

# FRONT — bp1.png (BOM extraction in action) — lower-left, hero
fw, fx, fy = 7.0, 0.9, 3.75
fh = fw * 0.428
rect(s, Inches(fx + 0.09), Inches(fy + 0.11),
     Inches(fw), Inches(fh), SHADOW_FRONT)
pic_front = s.shapes.add_picture(
    r"C:\Users\HP\Documents\PROJECTS\slide-desk\bp1.png",
    Inches(fx), Inches(fy), width=Inches(fw), height=Inches(fh))

footer_mark(s, "動作中", "BLUEPRINT · AT WORK")
page_mark(s, 10, TOTAL)
set_transition(s)

# ============================================================================
# 10 — Blueprint · now and next (operational)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "進捗 · 其ノ弐", "BLUEPRINT · NOW & NEXT", color=AI)
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "Blueprint — now and next.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=34, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
     "What's running, what's coming, who's doing it.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.3),
     "今、進めていること。次に取り掛かること。担うのは誰か。",
     size=11, color=FAINT, italic=True, spacing=100)

# Two columns: NOW · NEXT (team is on the next slide)
col_w = Inches(5.4); col_y = Inches(3.45)

def draw_bp_column(slide, x, header, header_color, items):
    shapes = []
    shapes.append(text(slide, x, col_y, col_w, Inches(0.32),
         header, size=10, bold=True, color=header_color, spacing=500))
    shapes.append(hairline(slide, x, col_y + Inches(0.42), Inches(0.55),
             color=header_color, thickness=Emu(19050)))
    iy = col_y + Inches(0.75)
    for en_line, jp_line in items:
        shapes.append(oval(slide, x, iy + Inches(0.1), Inches(0.09), header_color))
        shapes.append(text(slide, x + Inches(0.25), iy, Inches(5.0), Inches(0.4),
             en_line, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14,
             color=SUMI, bold=True, line_spacing=1.3))
        shapes.append(text(slide, x + Inches(0.25), iy + Inches(0.4), Inches(5.0),
             Inches(0.35), jp_line,
             size=11, color=FAINT, italic=True, spacing=50, line_spacing=1.3))
        iy += Inches(0.85)
    return shapes

x_now  = Inches(0.9)
x_next = Inches(0.9 + 5.4 + 0.7)

bp_now_shapes = draw_bp_column(s, x_now, "NOW · 進行中", GIN, [
    ("BP Check · BOM · Scan running in production",
     "BP Check / BOM / Scan が本番稼働中"),
    ("GPU / workflow decoupling — active rework",
     "GPU・ワークフロー分離の改修中"),
    ("Continuous testing + monitoring",
     "継続的なテストと監視"),
    ("Per-tenant configuration support",
     "テナント別設定への対応"),
])

bp_next_shapes = draw_bp_column(s, x_next, "NEXT · 60 DAYS · 次の60日", AI, [
    ("New workflows tailored for specific organisations",
     "特定組織向けの新ワークフロー追加"),
    ("Stabilise + harden recently shipped features",
     "新規機能の安定化と強化検証"),
    ("BP Check / BOM to first billed engineering teams",
     "BP Check・BOMを初の有償運用へ"),
])

footer_mark(s, "進捗 · BLUEPRINT", "NOW & NEXT")
page_mark(s, 11, TOTAL)
animate_on_click(s, [bp_now_shapes, bp_next_shapes])
set_transition(s)

# ============================================================================
# 12 — Blueprint · the team (Aita, Ambrose, Lauben, Otai)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "体制 · 其ノ弐", "BLUEPRINT · THE TEAM", color=AI)
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "The Blueprint team.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=34, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
     "Six contributors carrying the engine — four dedicated, two shared.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.3),
     "ブループリントを支える六人 — 専属四名、横断二名。",
     size=11, color=FAINT, italic=True, spacing=100)

# Six cards in a 2-row × 3-col grid
card_w = Inches(3.74); card_h = Inches(1.45); gap_x = 0.15; gap_y = 0.18
left_margin = 0.9
top_y = 3.35

def bp_card(row, col, kanji, name, role, hours, shared=False):
    x = Inches(left_margin + col * (3.74 + gap_x))
    y = Inches(top_y + row * (1.45 + gap_y))
    rect(s, x, y, card_w, card_h, KINARI_D)
    rect(s, x, y, card_w, Emu(28575), AI)
    # Name
    text(s, x + Inches(0.28), y + Inches(0.18),
         Inches(2.6), Inches(0.45),
         name, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=18,
         bold=True, color=SUMI)
    # Kanji top-right
    text(s, x + Inches(3.74 - 0.7), y + Inches(0.13),
         Inches(0.6), Inches(0.55),
         kanji, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22,
         color=AI, bold=True, align=PP_ALIGN.RIGHT)
    # Role (with optional shared marker)
    role_text = role + ("  *" if shared else "")
    text(s, x + Inches(0.28), y + Inches(0.62),
         Inches(3.74 - 0.55), Inches(0.3),
         role_text, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=11,
         italic=True, color=AI)
    # Hours
    text(s, x + Inches(0.28), y + Inches(0.96),
         Inches(3.74 - 0.55), Inches(0.4),
         hours, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15,
         bold=True, color=BRASS)

# Top row: dedicated specialist roles
bp_card(0, 0, "設計", "Aita",    "Backend & Engine Lead",   "32 h / week")
bp_card(0, 1, "土台", "Lauben",  "Infrastructure",          "40 h / week")
bp_card(0, 2, "補佐", "Otai",    "Backend Support",         "32 h / week")
# Bottom row: systems + the two shared contributors
bp_card(1, 0, "統率", "Ambrose", "Systems & Product",       "16 h / week")
bp_card(1, 1, "基盤", "Muzuva",  "Backend & ComfyUI",       "16 h / week", shared=True)
bp_card(1, 2, "表層", "Crispus", "Frontend",                "32 h / week", shared=True)

# Footnote about shared contributors
text(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.3),
     "*  Muzuva and Crispus also contribute to Apparel  ·  Apparelとも兼務",
     size=10, italic=True, color=FAINT, spacing=100)

# Bottom summary
hairline(s, Inches(0.9), Inches(6.78), Inches(11.5))
text(s, Inches(0.9), Inches(6.88), Inches(11.5), Inches(0.3),
     "6 contributors  ·  168 hrs / week  ·  Hold; +1 backend if client workflows scale.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=11, bold=True, color=AI)

footer_mark(s, "体制 · BLUEPRINT", "THE TEAM")
page_mark(s, 12, TOTAL)
set_transition(s)

# ============================================================================
# 11 — What we've learned (grounded engineering lessons)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI_D)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "学び", "WHAT WE'VE LEARNED", color=AI)
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.8),
     "Three lessons the quarter pressed into us.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=30, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.65), Inches(11), Inches(0.4),
     "今期、現場が教えてくれた三つのこと。",
     size=12, color=FAINT, italic=True, spacing=100)

lessons = [
    ("一", "Responsibility", "責任",
     "A process only one person can run is a risk the whole team carries.",
     "This quarter we retired two single-person dependencies by documenting them out loud.",
     "一人しか動かせない工程は、チーム全体のリスクでした。"),
    ("二", "Listening to the client", "顧客の声",
     "Customers find the problems our tests and monitors cannot.",
     "Our best fixes this quarter came from support tickets, not from dashboards.",
     "最良の修正は、監視画面ではなく顧客の声から生まれました。"),
    ("三", "Documentation is engineering", "文書化",
     "Docs ship with the code, or they never ship at all.",
     "Every shortcut we took on documentation became next sprint's delay.",
     "省いた文書は、次のスプリントの遅延になりました。"),
]
groups = []
y = Inches(3.25)
for kanji, title, jp_title, body, reason, jp in lessons:
    # Left column — kanji + title (width 3.8)
    k = text(s, Inches(0.9), y + Inches(0.05), Inches(0.9), Inches(0.9), kanji,
             latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=42, color=SHU, bold=True)
    t_ = text(s, Inches(1.85), y, Inches(2.9), Inches(0.45), title,
              latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=17, color=AI,
              bold=True, spacing=200)
    tj = text(s, Inches(1.85), y + Inches(0.5), Inches(2.9), Inches(0.4), jp_title,
              size=12, color=SHU, italic=True, spacing=200, bold=True)
    # Right column — body + reason + jp (width 7.6)
    lb = text(s, Inches(5.1), y + Inches(0.02), Inches(7.6), Inches(0.5), body,
              latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15, italic=True,
              color=SUMI, bold=True, line_spacing=1.35)
    lr = text(s, Inches(5.1), y + Inches(0.5), Inches(7.6), Inches(0.45), reason,
              latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, italic=True,
              color=FAINT, line_spacing=1.35)
    lj = text(s, Inches(5.1), y + Inches(0.92), Inches(7.6), Inches(0.35), jp,
              size=10, color=FAINT, spacing=100, italic=True)
    # Row separator (omit for last lesson to avoid overlap with footer)
    row_group = [k, t_, tj, lb, lr, lj]
    if kanji != "三":
        hr = hairline(s, Inches(0.9), y + Inches(1.15), Inches(11.5))
        row_group.append(hr)
    groups.append(row_group)
    y += Inches(1.22)

footer_mark(s, "学び", "WHAT WE'VE LEARNED")
page_mark(s, 13, TOTAL)
animate_on_click(s, groups)
set_transition(s)

# ============================================================================
# 12 — In Brief (circular journey visualisation)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "概要", "IN BRIEF")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "In Brief.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=38, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.55), Inches(11), Inches(0.4),
     "Four figures that map this quarter — the journey in numbers.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.9), Inches(11), Inches(0.3),
     "四つの数字でたどる、今期の歩み。",
     size=11, color=FAINT, italic=True, spacing=100)

# Journey path — thin horizontal line across slide, behind nodes
path_y_emu = Inches(4.85)
rect(s, Inches(1.9), path_y_emu, Inches(9.5), Emu(9525), HAIRLINE)

# Arrow-like small caret markers between circles to suggest forward motion
for cx in [3.75, 6.65, 9.55]:
    rect(s, Inches(cx - 0.05), Inches(4.78), Inches(0.1), Inches(0.15), SHU)

# 4 journey stations — node with big number, labels below
stations = [
    ("2",         "products live",       "in production",             "本番稼働中"),
    ("6",         "engineers",           "one end-to-end team",       "一つのチーム"),
    ("1 → 100+",  "one product photo",   "becomes a full campaign",   "一枚から百枚超へ"),
    ("hrs → min", "blueprint reviews",   "hours down to minutes",     "数時間から数分へ"),
]
start_x = 2.3
step = 2.9
node_size = 1.4
for i, (big, line1, line2, jp) in enumerate(stations):
    cx = start_x + i * step
    cy = 4.85
    # Bordered circle
    nd = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(cx - node_size/2), Inches(cy - node_size/2),
                            Inches(node_size), Inches(node_size))
    nd.fill.solid(); nd.fill.fore_color.rgb = KINARI_D
    nd.line.color.rgb = SHU
    nd.line.width = Pt(1.25)

    # Small chapter numeral above circle
    chapter = ["一", "二", "三", "四"][i]
    text(s, Inches(cx - 0.5), Inches(cy - 1.35), Inches(1.0), Inches(0.35),
         chapter, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=16,
         color=BRASS, bold=True, align=PP_ALIGN.CENTER, spacing=300)

    # Big number inside circle (font scaled down for longer strings)
    num_size = 30 if len(big) <= 2 else 22 if len(big) <= 5 else 18
    text(s, Inches(cx - 0.75), Inches(cy - 0.45), Inches(1.5), Inches(0.9),
         big, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=num_size, bold=True,
         color=AI, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Labels below
    text(s, Inches(cx - 1.35), Inches(cy + 0.85), Inches(2.7), Inches(0.35),
         line1, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13, bold=True,
         color=SUMI, align=PP_ALIGN.CENTER)
    text(s, Inches(cx - 1.35), Inches(cy + 1.2), Inches(2.7), Inches(0.35),
         line2, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=11, italic=True,
         color=FAINT, align=PP_ALIGN.CENTER)
    text(s, Inches(cx - 1.35), Inches(cy + 1.58), Inches(2.7), Inches(0.3),
         jp, size=10, color=FAINT, italic=True,
         align=PP_ALIGN.CENTER, spacing=100)

# Footer note
text(s, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.3),
     "Figures directional for this note; final numbers confirmed at quarter-end.   数字は目安、四半期締めに確定。",
     size=10, color=GIN, italic=True, spacing=100, align=PP_ALIGN.CENTER)

footer_mark(s, "概要", "IN BRIEF")
page_mark(s, 14, TOTAL)
set_transition(s)

# ============================================================================
# 15 — The team picture (allocation across both products)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "配置", "THE TEAM PICTURE")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.7),
     "Six people, two products, two crossing both.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=30, bold=True, color=SUMI)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
     "Where the team currently sits — and where capacity could shift.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.3),
     "二つのプロダクトと、それを横断する二人。現状と今後の調整余地。",
     size=11, color=FAINT, italic=True, spacing=100)

# Two team columns
col_top = Inches(3.4)
col_w = Inches(5.5)
col_a_x = Inches(0.9)
col_b_x = Inches(0.9 + 5.5 + 0.5)  # 6.9

# APPAREL column header
text(s, col_a_x, col_top, col_w, Inches(0.4),
     "APPAREL  ·  2 人  ·  24 hrs / week",
     size=11, bold=True, color=SHU, spacing=400)
hairline(s, col_a_x, col_top + Inches(0.5), Inches(0.6),
         color=SHU, thickness=Emu(19050))

apparel_team = [
    ("Muzuva",  "Backend & ComfyUI Lead",  "16h", True),
    ("Crispus", "Frontend",                "8h",  True),
]
ay = col_top + Inches(0.78)
for name, role, hrs, shared in apparel_team:
    oval(s, col_a_x, ay + Inches(0.1), Inches(0.09), SHU)
    text(s, col_a_x + Inches(0.25), ay, Inches(1.5), Inches(0.35),
         name, latin=LATIN_DISPLAY, ja=JA_DISPLAY,
         size=14, bold=True, color=SUMI)
    text(s, col_a_x + Inches(1.7), ay, Inches(3.0), Inches(0.35),
         role + (" *" if shared else ""),
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, color=FAINT)
    text(s, col_a_x + Inches(4.6), ay, Inches(0.8), Inches(0.35),
         hrs, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13,
         color=BRASS, bold=True, align=PP_ALIGN.RIGHT)
    ay += Inches(0.45)

# BLUEPRINT column header
text(s, col_b_x, col_top, col_w, Inches(0.4),
     "BLUEPRINT  ·  6 人  ·  168 hrs / week",
     size=11, bold=True, color=AI, spacing=400)
hairline(s, col_b_x, col_top + Inches(0.5), Inches(0.6),
         color=AI, thickness=Emu(19050))

blueprint_team = [
    ("Aita",    "Backend & Engine Lead",  "32h", False),
    ("Lauben",  "Infrastructure",         "40h", False),
    ("Otai",    "Backend Support",        "32h", False),
    ("Ambrose", "Systems & Product",      "16h", False),
    ("Crispus", "Frontend",               "32h", True),
    ("Muzuva",  "Backend & ComfyUI",      "16h", True),
]
by = col_top + Inches(0.78)
for name, role, hrs, shared in blueprint_team:
    oval(s, col_b_x, by + Inches(0.1), Inches(0.09), AI)
    text(s, col_b_x + Inches(0.25), by, Inches(1.5), Inches(0.35),
         name, latin=LATIN_DISPLAY, ja=JA_DISPLAY,
         size=14, bold=True, color=SUMI)
    text(s, col_b_x + Inches(1.7), by, Inches(3.0), Inches(0.35),
         role + (" *" if shared else ""),
         latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, color=FAINT)
    text(s, col_b_x + Inches(4.6), by, Inches(0.8), Inches(0.35),
         hrs, latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13,
         color=BRASS, bold=True, align=PP_ALIGN.RIGHT)
    by += Inches(0.45)

# Note about shared contributors
text(s, Inches(0.9), Inches(6.18), Inches(11.5), Inches(0.3),
     "*  Muzuva and Crispus split time across both products.   Unique headcount: 6.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=11, italic=True, color=FAINT)

# Open question for the JP/UG conversation
hairline(s, Inches(0.9), Inches(6.5), Inches(11.5))
text(s, Inches(0.9), Inches(6.62), Inches(11.5), Inches(0.35),
     "Open question  ·  Where does Japan-side capacity best slot in over the next 60 days?",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=12, bold=True, color=AI)
text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.3),
     "今後60日、日本側のリソースをどこに配置するのが最善か。",
     size=10, color=FAINT, italic=True, spacing=100)

footer_mark(s, "配置", "TEAM PICTURE")
page_mark(s, 15, TOTAL)
set_transition(s)

# ============================================================================
# 16 — The next 60 days (concrete)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, AI)
rect(s, 0, 0, Inches(0.4), SLIDE_H, SHU)

kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "これから", "THE NEXT 60 DAYS", color=KINARI)
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5), color=AI_SOFT)

text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.8),
     "Three concrete moves in the next 60 days.",
     latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=30, bold=True, color=KINARI)
text(s, Inches(0.9), Inches(2.65), Inches(11), Inches(0.4),
     "今後60日間で、確実に進める三手。",
     size=12, color=HAIRLINE, italic=True, spacing=100)

moves = [
    ("Apparel",
     "Users are already on board.  We sharpen the background engine —",
     "new perspectives, richer scenes, a visibly better experience.",
     "背景生成の精度と構図を引き上げ、体験をもう一段高めます。"),
    ("Blueprint",
     "New workflows tailored for specific organisations.",
     "Stabilise, test, and harden what's already shipped.",
     "特定組織向けの新ワークフロー追加と、既存機能の安定化・検証。"),
    ("The team",
     "Close the last single-person dependencies.",
     "Turn the learnings from both products into shared team playbooks.",
     "単独依存を解消し、両プロダクトの学びをチームの手引きに。"),
]
groups = []
y = Inches(3.3)
for title, a, b, jp in moves:
    t_ = text(s, Inches(0.9), y, Inches(3.0), Inches(0.5), title,
              latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=22, bold=True,
              color=SHU)
    la = text(s, Inches(3.9), y + Inches(0.05), Inches(9.2), Inches(0.45), a,
              latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=15, bold=True,
              color=KINARI)
    lb = text(s, Inches(3.9), y + Inches(0.45), Inches(9.2), Inches(0.45), b,
              latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True,
              color=HAIRLINE)
    lj = text(s, Inches(3.9), y + Inches(0.9), Inches(9.2), Inches(0.35), jp,
              size=10, color=HAIRLINE, italic=True, spacing=100)
    h = hairline(s, Inches(0.9), y + Inches(1.2), Inches(12.2),
                 color=AI_SOFT, thickness=Emu(6350))
    groups.append([t_, la, lb, lj, h])
    y += Inches(1.15)

footer_mark(s, "これから", "THE NEXT 60 DAYS", color=HAIRLINE)
page_mark(s, 16, TOTAL, color=HAIRLINE)
animate_on_click(s, groups)
set_transition(s)

# ============================================================================
# 17 — One request (closing)
# ============================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, KINARI)
kicker(s, Inches(0.9), Inches(0.7), Inches(6),
       "お伺い", "ONE REQUEST")
hairline(s, Inches(0.9), Inches(1.45), Inches(11.5))

q1 = text(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.85),
          "Which of the two would you like us to open next?",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=36, bold=True, color=SUMI,
          align=PP_ALIGN.CENTER)
q2 = text(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(0.5),
          "どちらの作品を、次にご覧になりたいですか。",
          size=15, color=FAINT, italic=True, align=PP_ALIGN.CENTER, spacing=100)
q3 = text(s, Inches(0.9), Inches(3.65), Inches(11.5), Inches(0.45),
          "We'd be glad to walk the team deeper into whichever one interests you.",
          latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=13, italic=True,
          color=FAINT, align=PP_ALIGN.CENTER)

tile_w = Inches(5.0); gap = Inches(0.4)
tx1 = Inches((13.333 - 2*5.0 - 0.4) / 2)
tx2 = tx1 + tile_w + gap
ty = Inches(4.25); tile_h = Inches(2.35)

# Tile 1 — Apparel
t1a = rect(s, tx1, ty, tile_w, tile_h, KINARI_D)
t1b = rect(s, tx1, ty, tile_w, Emu(22860), SHU)
t1c = text(s, tx1 + Inches(0.4), ty + Inches(0.25), tile_w, Inches(0.4),
           "其ノ壱   ·   APPAREL", size=10, color=SHU, spacing=600, bold=True)
t1d = text(s, tx1 + Inches(0.4), ty + Inches(0.65), tile_w, Inches(1.0),
           "Apparel",
           latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=38, bold=True, color=SUMI)
t1e = text(s, tx1 + Inches(0.4), ty + Inches(1.5), tile_w, Inches(0.45),
           "One product photo, a full campaign.",
           latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
t1f = text(s, tx1 + Inches(0.4), ty + Inches(1.87), tile_w, Inches(0.4),
           "一枚の写真が、キャンペーンに。",
           size=11, color=FAINT, italic=True, spacing=100)

# Tile 2 — Blueprint
t2a = rect(s, tx2, ty, tile_w, tile_h, KINARI_D)
t2b = rect(s, tx2, ty, tile_w, Emu(22860), AI)
t2c = text(s, tx2 + Inches(0.4), ty + Inches(0.25), tile_w, Inches(0.4),
           "其ノ弐   ·   BLUEPRINT", size=10, color=AI, spacing=600, bold=True)
t2d = text(s, tx2 + Inches(0.4), ty + Inches(0.65), tile_w, Inches(1.0),
           "Blueprint",
           latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=38, bold=True, color=SUMI)
t2e = text(s, tx2 + Inches(0.4), ty + Inches(1.5), tile_w, Inches(0.45),
           "One blueprint, a parts list and inspection.",
           latin=LATIN_DISPLAY, ja=JA_DISPLAY, size=14, italic=True, color=FAINT)
t2f = text(s, tx2 + Inches(0.4), ty + Inches(1.87), tile_w, Inches(0.4),
           "一枚の図面が、部品表と検査に。",
           size=11, color=FAINT, italic=True, spacing=100)

hanko(s, Inches(12.05), Inches(6.25), size=Inches(0.7), char="承")

text(s, Inches(0.9), Inches(6.9), Inches(10), Inches(0.3),
     "alanda.ambrose@aibos.co.jp   ·   Create Project",
     size=10, color=GIN, spacing=300)

page_mark(s, 17, TOTAL, color=GIN)
animate_on_click(s,
    [[q1, q2, q3],
     [t1a, t1b, t1c, t1d, t1e, t1f],
     [t2a, t2b, t2c, t2d, t2e, t2f]])
set_transition(s)

# --- Save --------------------------------------------------------------------
out = r"C:\Users\HP\Documents\PROJECTS\slide-desk\create_project_note_v15.pptx"
prs.save(out)
print(f"Saved: {out}")
